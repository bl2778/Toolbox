"""PPTX parsing utilities for the WR module."""

from __future__ import annotations

import re
import uuid
from typing import List, Dict

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from .models import SlideElement, SlimSlide

WORD_RE = re.compile(r"[A-Za-z']+")
NOTE_RE = re.compile(r"^note\s*:", re.IGNORECASE)
URL_RE = re.compile(r"https?://|www\.", re.IGNORECASE)
EXPORT_ARTIFACT_RE = re.compile(r"[A-Za-z]+_\d{1,}")


def _generate_id() -> str:
    return uuid.uuid4().hex[:8]


def _split_lines(text: str) -> List[str]:
    if not text:
        return []
    cleaned = text.replace("\u000b", "\n")
    return [line.strip() for line in cleaned.splitlines() if line.strip()]


def _is_metrics_only(text: str) -> bool:
    # Drop strings that are primarily numeric/metric content
    cleaned = re.sub(r"[\s%$,.;:()\-]+", " ", text).strip()
    if not cleaned:
        return True
    # If no alphabetic characters remain, treat as metrics only
    return not re.search(r"[A-Za-z]", cleaned)


def _is_table_header(text: str) -> bool:
    lowered = text.lower()
    keywords = [
        "time to value",
        "size of prize",
        "cost bucket",
        "owner",
        "status",
        "timeline",
        "priority",
    ]
    return any(key in lowered for key in keywords) and len(WORD_RE.findall(text)) <= 6


def _is_candidate_text(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    if URL_RE.search(stripped):
        return False
    if NOTE_RE.match(stripped) and len(WORD_RE.findall(stripped)) < 10:
        return False
    if EXPORT_ARTIFACT_RE.search(stripped):
        return False
    words = WORD_RE.findall(stripped)
    if len(words) < 5:
        return False
    if _is_metrics_only(stripped):
        return False
    if _is_table_header(stripped):
        return False
    return True


def _detect_type(shape) -> str:
    if getattr(shape, "is_placeholder", False):
        try:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type in {
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
                PP_PLACEHOLDER.SUBTITLE,
                PP_PLACEHOLDER.VERTICAL_TITLE,
            }:
                return "Title/Subtitle"
        except Exception:
            pass
    return "Body"


def _normalize_table_line(line: str) -> str:
    idx = line.find(":")
    if idx >= 0:
        return line[idx + 1 :].strip()
    return line.strip()


def extract_slim_json(ppt_path: str) -> List[Dict]:
    presentation = Presentation(ppt_path)
    slides: List[SlimSlide] = []

    for index, slide in enumerate(presentation.slides, start=1):
        elements: List[SlideElement] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame and shape.text_frame.text:
                lines = _split_lines(shape.text_frame.text)
                filtered = [line for line in lines if _is_candidate_text(line)]
                if filtered:
                    elements.append(
                        SlideElement(
                            id=_generate_id(),
                            type=_detect_type(shape),
                            text="\n".join(filtered),
                        )
                    )
            if getattr(shape, "has_table", False) and shape.has_table:
                table_lines: List[str] = []
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if not cell_text:
                            continue
                        normalized = _normalize_table_line(cell_text)
                        if _is_candidate_text(normalized):
                            table_lines.append(f"Row ?, Col ?: {normalized}")
                if table_lines:
                    elements.append(
                        SlideElement(
                            id=_generate_id(),
                            type="Table",
                            text="\n".join(table_lines),
                        )
                    )
        if elements:
            slides.append(SlimSlide(slide_number=index, elements=elements))

    return [
        {
            "slide_number": slide.slide_number,
            "elements": [
                {"id": element.id, "type": element.type, "text": element.text}
                for element in slide.elements
            ],
        }
        for slide in slides
    ]
