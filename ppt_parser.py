"""
PowerPoint Text Extractor for ZD Tool
-------------------------------------
Extracts structured text from PowerPoint files for Zero-Defect checking.
Supports conversion to JSON format optimized for AI analysis.
"""

import json
import os
import re
from pathlib import Path
from typing import List, Dict, Optional, Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.shapes import PP_PLACEHOLDER
except ImportError:
    raise ImportError("python-pptx is not installed. Run: pip install python-pptx")


class PPTExtractor:
    def __init__(self):
        pass

    def extract_text_recursive(self, shape, depth: int = 0, max_depth: int = 10) -> List[Dict[str, str]]:
        """Recursively extract text from a shape (handling groups, tables, charts)."""
        chunks = []

        # Prevent infinite recursion
        if depth > max_depth:
            print(f"[WARNING] Max recursion depth reached for shape {shape.shape_id}")
            return chunks

        sid = f"Shape ID {shape.shape_id}"

        try:
            # 1) Plain text frames
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text_frame.text.strip():
                text = shape.text_frame.text.strip()
                type_hint = "Body"
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    try:
                        ph = shape.placeholder_format
                        if ph.type in {
                            PP_PLACEHOLDER.TITLE,
                            PP_PLACEHOLDER.CENTER_TITLE,
                            PP_PLACEHOLDER.SUBTITLE,
                            PP_PLACEHOLDER.VERTICAL_TITLE,
                        }:
                            type_hint = "Title/Subtitle"
                        elif ph.type == PP_PLACEHOLDER.BODY:
                            type_hint = "Body Placeholder"
                        elif ph.type == PP_PLACEHOLDER.OBJECT and hasattr(shape, 'name') and "Title" in shape.name:
                            type_hint = "Object Title"
                    except Exception:
                        pass  # Placeholder format access can fail

                chunks.append({"id": sid, "type": type_hint, "text": text})

            # 2) Table cells (with limits to prevent hanging)
            elif hasattr(shape, 'has_table') and shape.has_table:
                tbl_txt = []
                max_rows = min(50, len(shape.table.rows))  # Limit to 50 rows
                for r in range(max_rows):
                    row = shape.table.rows[r]
                    max_cols = min(20, len(row.cells))  # Limit to 20 columns
                    for c in range(max_cols):
                        try:
                            cell = row.cells[c]
                            cell_txt = cell.text_frame.text.strip()
                            if cell_txt:
                                tbl_txt.append(f"Row {r+1}, Col {c+1}: {cell_txt}")
                        except Exception as e:
                            print(f"[WARNING] Error reading table cell [{r},{c}]: {e}")
                            continue
                if tbl_txt:
                    chunks.append({"id": sid, "type": "Table", "text": "\\n".join(tbl_txt)})

            # 3) Grouped shapes (with depth tracking)
            elif hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                if hasattr(shape, 'shapes'):
                    max_shapes = min(100, len(shape.shapes))  # Limit to 100 shapes per group
                    for i in range(max_shapes):
                        try:
                            s = shape.shapes[i]
                            chunks.extend(self.extract_text_recursive(s, depth + 1, max_depth))
                        except Exception as e:
                            print(f"[WARNING] Error processing grouped shape {i}: {e}")
                            continue

            # 4) Chart title (limited)
            elif hasattr(shape, 'has_chart') and shape.has_chart:
                try:
                    ch = shape.chart
                    if hasattr(ch, 'has_title') and ch.has_title and ch.chart_title.text_frame.text.strip():
                        chunks.append(
                            {
                                "id": sid,
                                "type": "Chart Info",
                                "text": f"Chart Title: {ch.chart_title.text_frame.text.strip()}",
                            }
                        )
                except Exception as e:
                    print(f"[WARNING] Error processing chart: {e}")

        except Exception as e:
            print(f"[ERROR] Error processing shape {sid}: {e}")

        return chunks

    def extract_powerpoint_text(self, pptx_path: str) -> Optional[List[Dict[str, Any]]]:
        """Extract text from PowerPoint file and return structured data."""
        import time

        if not Path(pptx_path).exists():
            raise FileNotFoundError(f"File not found: {pptx_path}")

        print(f"[INFO] Starting PowerPoint extraction: {pptx_path}")
        start_time = time.time()

        try:
            prs = Presentation(pptx_path)
            print(f"[INFO] Successfully opened presentation with {len(prs.slides)} slides")
        except Exception as exc:
            raise ValueError(f"Could not open presentation: {exc}")

        slides_data = []
        max_slides = min(500, len(prs.slides))  # Limit to 500 slides maximum

        for idx in range(max_slides):
            slide = prs.slides[idx]
            slide_info = {"slide_number": idx + 1, "elements": [], "notes": None}

            # Check timeout (5 minutes max for parsing)
            if time.time() - start_time > 300:
                print(f"[WARNING] Extraction timeout after 5 minutes, stopping at slide {idx + 1}")
                break

            try:
                # Extract from shapes with limits
                max_shapes = min(200, len(slide.shapes))  # Limit to 200 shapes per slide
                for shape_idx in range(max_shapes):
                    try:
                        shp = slide.shapes[shape_idx]
                        extracted = self.extract_text_recursive(shp, depth=0, max_depth=5)
                        slide_info["elements"].extend(extracted)
                    except Exception as e:
                        print(f"[WARNING] Error processing shape {shape_idx} on slide {idx + 1}: {e}")
                        continue

                # Extract notes with timeout protection
                if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
                    try:
                        notes_tf = slide.notes_slide.notes_text_frame
                        if notes_tf and notes_tf.text.strip():
                            slide_info["notes"] = notes_tf.text.strip()
                    except Exception as e:
                        print(f"[WARNING] Error extracting notes from slide {idx + 1}: {e}")

            except Exception as e:
                print(f"[WARNING] Error processing slide {idx + 1}: {e}")

            # Always add slide even if empty (for consistent page numbering)
            slides_data.append(slide_info)

        elapsed = time.time() - start_time
        print(f"[INFO] PowerPoint extraction completed in {elapsed:.2f} seconds, processed {len(slides_data)} slides")

        return slides_data

    def convert_to_zd_format(self, slides_data: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Convert extracted data to ZD Tool format."""
        zd_slides = []

        for slide in slides_data:
            slide_number = slide["slide_number"]

            # Separate title/tagline from body content
            titles = []
            body_parts = []

            for element in slide["elements"]:
                if element["type"] == "Title/Subtitle":
                    titles.append(element["text"])
                else:
                    # Filter out purely alphanumeric markers like "A.", "I-1"
                    text = element["text"]
                    if not re.match(r'^[A-Z]\.?$|^[IVX]+-?\d*\.?$', text.strip()):
                        body_parts.append(text)

            # Combine titles as tagline
            tagline = " | ".join(titles) if titles else ""

            # Combine body content
            body_other = "\\n".join(body_parts) if body_parts else ""

            # Notes (marked as "Do not review" per spec)
            speaker_notes = "Do not review" if slide["notes"] else ""

            zd_slide = {
                "page_number": slide_number,
                "tagline": tagline,
                "body_other": body_other,
                "speaker_notes": speaker_notes
            }

            zd_slides.append(zd_slide)

        return zd_slides

    def get_word_count(self, text: str, language: str = "english") -> int:
        """Count words in English text or characters in Chinese text."""
        if not text or not text.strip():
            return 0

        if language == "chinese":
            # For Chinese, count characters (excluding spaces and punctuation)
            import re
            # Remove all whitespace and common punctuation
            chinese_chars = re.sub(r'[\s\.\,\!\?\:\;\"\'\(\)\[\]\{\}]', '', text)
            return len(chinese_chars)
        else:
            # For English, count words
            return len(text.split())

    def get_slide_stats(self, zd_slides: List[Dict[str, Any]], language: str = "english") -> Dict[str, Any]:
        """Get statistics about the slides."""
        total_slides = len(zd_slides)
        total_words = 0

        for slide in zd_slides:
            slide_words = (self.get_word_count(slide["tagline"], language) +
                          self.get_word_count(slide["body_other"], language))
            total_words += slide_words

        avg_words_per_slide = total_words / total_slides if total_slides > 0 else 0

        # Determine the unit based on language
        unit_name = "characters" if language == "chinese" else "words"

        return {
            "total_slides": total_slides,
            "total_words": total_words,
            "avg_words_per_slide": round(avg_words_per_slide, 1),
            "recommended_mode": "fast" if avg_words_per_slide < 100 else "precise",
            "unit": unit_name
        }


class TextChunker:
    """Handles text chunking for ZD analysis."""

    def __init__(self, language: str = "english"):
        self.language = language

        # Chunking configuration
        self.fast_config = {
            "target_words_min": 4000,
            "target_words_max": 6500,
            "max_pages": 10,
            "overlap_pages": 1
        }

        self.precise_config = {
            "target_words_min": 2500,
            "target_words_max": 4000,
            "max_pages": 5,
            "overlap_pages": 1
        }

    def count_slide_words(self, slide: Dict[str, Any]) -> int:
        """Count words/characters in tagline and body_other only."""
        if self.language == "chinese":
            import re
            # For Chinese, count characters (excluding spaces and punctuation)
            tagline_text = slide.get("tagline", "")
            body_text = slide.get("body_other", "")
            combined_text = tagline_text + body_text
            # Remove all whitespace and common punctuation
            chinese_chars = re.sub(r'[\s\.\,\!\?\:\;\"\'\(\)\[\]\{\}]', '', combined_text)
            return len(chinese_chars)
        else:
            # For English, count words
            tagline_words = len(slide["tagline"].split()) if slide["tagline"] else 0
            body_words = len(slide["body_other"].split()) if slide["body_other"] else 0
            return tagline_words + body_words

    def create_chunks(self, zd_slides: List[Dict[str, Any]], mode: str = "fast") -> List[Dict[str, Any]]:
        """Create chunks based on the specified mode."""
        config = self.fast_config if mode == "fast" else self.precise_config
        chunks = []
        chunk_id = 1

        i = 0
        while i < len(zd_slides):
            chunk_slides = []
            chunk_words = 0
            chunk_pages = 0

            # Add slides to chunk until we hit limits
            while (i < len(zd_slides) and
                   chunk_pages < config["max_pages"] and
                   (chunk_words < config["target_words_min"] or
                    (chunk_words + self.count_slide_words(zd_slides[i]) <= config["target_words_max"]))):

                slide_words = self.count_slide_words(zd_slides[i])

                # Avoid infinite loop: if chunk is empty and adding this slide exceeds max, add it anyway
                if not chunk_slides and chunk_words + slide_words > config["target_words_max"]:
                    chunk_slides.append(zd_slides[i])
                    chunk_words += slide_words
                    chunk_pages += 1
                    i += 1
                    break
                elif chunk_words + slide_words <= config["target_words_max"]:
                    chunk_slides.append(zd_slides[i])
                    chunk_words += slide_words
                    chunk_pages += 1
                    i += 1
                else:
                    break

            if chunk_slides:
                chunk = {
                    "chunk_id": f"ck_{chunk_id:04d}",
                    "mode": mode,
                    "page_start": chunk_slides[0]["page_number"],
                    "page_end": chunk_slides[-1]["page_number"],
                    "page_numbers": [slide["page_number"] for slide in chunk_slides],
                    "word_count": chunk_words,
                    "slides": chunk_slides
                }
                chunks.append(chunk)
                chunk_id += 1

                # Handle overlap: move back by overlap_pages - 1, but avoid creating redundant single-page chunks
                overlap = min(config["overlap_pages"], len(chunk_slides) - 1)
                # Only apply overlap if we haven't reached the end of all slides
                if i < len(zd_slides):
                    # Check if moving back would create a redundant chunk with only the last page(s)
                    remaining_slides = len(zd_slides) - i
                    if remaining_slides > overlap:
                        i -= overlap
                    # If remaining slides <= overlap, don't move back to avoid redundant chunks

        return chunks


def extract_ppt_for_zd(file_path: str, mode: str = "fast", language: str = "english") -> Dict[str, Any]:
    """Main function to extract and chunk PPT for ZD analysis."""
    import time

    start_time = time.time()
    unit_name = "characters" if language == "chinese" else "words"
    print(f"[INFO] Starting ZD extraction for: {file_path} (mode: {mode}, language: {language})")

    try:
        extractor = PPTExtractor()
        chunker = TextChunker(language=language)

        # Extract raw data with timeout
        print("[INFO] Step 1: Extracting PowerPoint text...")
        raw_slides = extractor.extract_powerpoint_text(file_path)
        if not raw_slides:
            raise ValueError("No extractable text found in the presentation")

        print(f"[INFO] Step 2: Converting {len(raw_slides)} slides to ZD format...")
        # Convert to ZD format
        zd_slides = extractor.convert_to_zd_format(raw_slides)

        print("[INFO] Step 3: Calculating statistics...")
        # Get statistics
        stats = extractor.get_slide_stats(zd_slides, language=language)

        print(f"[INFO] Step 4: Creating chunks in {mode} mode...")
        # Create chunks with timeout check
        if time.time() - start_time > 300:  # 5 minute total timeout
            raise TimeoutError("Total extraction time exceeded 5 minutes")

        chunks = chunker.create_chunks(zd_slides, mode)

        elapsed = time.time() - start_time
        print(f"[INFO] ZD extraction completed successfully in {elapsed:.2f} seconds")
        print(f"[INFO] Results: {len(zd_slides)} slides, {len(chunks)} chunks, {stats.get('total_words', 0)} {unit_name}")

        return {
            "success": True,
            "stats": stats,
            "slides": zd_slides,
            "chunks": chunks,
            "total_chunks": len(chunks)
        }

    except Exception as e:
        elapsed = time.time() - start_time
        print(f"[ERROR] ZD extraction failed after {elapsed:.2f} seconds: {str(e)}")
        return {
            "success": False,
            "error": str(e),
            "stats": None,
            "slides": [],
            "chunks": [],
            "total_chunks": 0
        }